# -*- coding: utf-8 -*-
"""PowerPoint 週報產生器(python-pptx,原生圖表)。

結構:封面 → 總體結論 → 產業聲量圖 → 競品 SOV(各組一頁)→ 負面警示。
圖表為 PowerPoint 原生圖表(非圖片),開啟端以本機字型渲染,無伺服器字型問題。
"""
from collections import Counter
from datetime import datetime, timezone, timedelta

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from config import REPORT_SOV_GROUPS
from analysis.companies import make_matcher
from analysis.summarizer import detect_themes
from analysis.alerts import negative_alerts

TZ_TAIPEI = timezone(timedelta(hours=8))


def _bullet_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
    return slide


def build_ppt(articles, start_date, end_date, watchlist=None, out_path=None,
              articles_prev=None):
    prs = Presentation()
    now = datetime.now(TZ_TAIPEI)

    # 封面
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "台灣產業趨勢監測週報"
    cover.placeholders[1].text = (f"報告期間:{start_date} ~ {end_date}\n"
                                  f"產出:{now.strftime('%Y-%m-%d %H:%M')}(台北時間)\n"
                                  "資料來源:公開新聞媒體與官方開放資料")

    classified = [a for a in articles if a.get("industry")]
    ind_counts = Counter(a["industry"] for a in classified)
    themes = detect_themes(articles)

    # 總體結論
    lines = [f"本期共蒐集 {len(articles)} 篇公開報導,可歸類 {len(classified)} 篇"]
    for ind, cnt in ind_counts.most_common(3):
        lines.append(f"{ind}:{cnt} 篇"
                     f"(佔可歸類 {cnt / max(len(classified), 1) * 100:.1f}%)")
    if themes:
        lines.append("最熱主題:" + "、".join(f"{t[0]}({t[1]} 篇)" for t in themes[:3]))
    lines.append("※ 本頁由系統統計自動生成,數據可溯源至報導原文")
    _bullet_slide(prs, "總體結論", lines)

    # 產業聲量(原生橫條圖)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "各產業報導篇數"
    ranked = ind_counts.most_common()
    data = CategoryChartData()
    data.categories = [k for k, _ in reversed(ranked)]
    data.add_series("篇數", [v for _, v in reversed(ranked)])
    slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                           Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.4), data)

    # 競品 SOV(各組一頁)
    def find_aliases(name):
        for comps in (watchlist or {}).values():
            if name in comps:
                return comps[name]
        return [name]

    def count(arts, match):
        return sum(1 for a in arts
                   if match(f"{a.get('title', '')} {a.get('summary', '')} {a.get('source', '')}"))

    for group_name, names in REPORT_SOV_GROUPS.items():
        rows = []
        for n in names:
            m = make_matcher(find_aliases(n))
            rows.append((n, count(articles, m),
                         count(articles_prev, m) if articles_prev else None))
        rows.sort(key=lambda x: -x[1])
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        gtotal = sum(c for _, c, _ in rows)
        slide.shapes.title.text = f"競品聲量對比(SOV)— {group_name}(合計 {gtotal} 篇)"
        data = CategoryChartData()
        data.categories = [n for n, _, _ in rows]
        data.add_series("本期篇數", [c for _, c, _ in rows])
        if articles_prev:
            data.add_series("前期篇數", [p or 0 for _, _, p in rows])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.2), data).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    # 負面警示
    alerts = negative_alerts(articles)
    if alerts:
        lines = [f"{al['company']}:疑似負面 {al['count']} 篇"
                 f"(命中:{'、'.join(al['keywords'][:4])})" for al in alerts[:8]]
        lines.append("※ 規則式偵測僅供初步示警,請人工確認報導立場")
    else:
        lines = ["本期未偵測到觀察名單公司之負面訊號"]
    _bullet_slide(prs, "負面聲量警示", lines)

    if out_path is None:
        out_path = f"台灣產業趨勢監測週報_{start_date}_{end_date}.pptx"
    prs.save(out_path)
    return str(out_path)
